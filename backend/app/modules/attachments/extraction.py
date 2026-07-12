from __future__ import annotations

import io
import zipfile
from pathlib import Path
from app.modules.attachments.domain import ExtractedContent
from app.shared.config import AttachmentSettings
from app.shared.exceptions import NonRetryableExecutionError


OOXML_MIME = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
IMAGE_MIME = {"JPEG": "image/jpeg", "PNG": "image/png", "WEBP": "image/webp"}


class SafeAttachmentExtractor:
    parser_version = "mvp-python-1"

    def __init__(self, settings: AttachmentSettings) -> None:
        self.settings = settings

    def inspect(self, *, file_name: str, data: bytes) -> str:
        extension = Path(file_name).suffix.lower()
        if extension not in self.settings.allowed_extensions:
            raise _rejected("unsupported_extension")
        if not data or len(data) > self.settings.max_file_bytes:
            raise _rejected("file_size_exceeded")
        if extension in OOXML_MIME:
            self._assert_safe_ooxml(extension, data)
            return OOXML_MIME[extension]
        if extension in {".md", ".markdown"}:
            data.decode("utf-8")
            return "text/markdown"
        if extension in {".jpg", ".jpeg", ".png", ".webp"}:
            return self._inspect_image(data)
        raise _rejected("unsupported_extension")

    def extract(self, *, file_name: str, data: bytes) -> ExtractedContent:
        extension = Path(file_name).suffix.lower()
        self.inspect(file_name=file_name, data=data)
        if extension in {".jpg", ".jpeg", ".png", ".webp"}:
            raise _rejected("image_not_interpreted")
        if extension in {".md", ".markdown"}:
            return self._bounded(data.decode("utf-8"), [{"type": "markdown"}])
        if extension == ".docx":
            return self._extract_docx(data)
        if extension == ".xlsx":
            return self._extract_xlsx(data)
        if extension == ".pptx":
            return self._extract_pptx(data)
        raise _rejected("unsupported_extension")

    def normalize_image(self, *, data: bytes) -> tuple[bytes, str]:
        try:
            from PIL import Image
        except ModuleNotFoundError as exc:
            raise RuntimeError("Pillow is required for image validation") from exc
        with Image.open(io.BytesIO(data)) as image:
            image.verify()
        with Image.open(io.BytesIO(data)) as image:
            if image.width * image.height > self.settings.max_image_pixels:
                raise _rejected("image_pixels_exceeded")
            image_format = str(image.format or "").upper()
            mime = IMAGE_MIME.get(image_format)
            if mime is None:
                raise _rejected("unsupported_image_format")
            clean = image.copy()
            output = io.BytesIO()
            save_format = "JPEG" if image_format == "JPEG" else image_format
            if save_format == "JPEG" and clean.mode not in {"RGB", "L"}:
                clean = clean.convert("RGB")
            clean.save(output, format=save_format)
            return output.getvalue(), mime

    def _assert_safe_ooxml(self, extension: str, data: bytes) -> None:
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as archive:
                total = sum(item.file_size for item in archive.infolist())
                names = {item.filename.lower() for item in archive.infolist()}
        except zipfile.BadZipFile as exc:
            raise _rejected("invalid_ooxml") from exc
        if total > self.settings.max_uncompressed_bytes:
            raise _rejected("uncompressed_size_exceeded")
        if any("vbaproject" in name or "oleobject" in name for name in names):
            raise _rejected("active_content_rejected")
        required = {
            ".docx": "word/document.xml",
            ".xlsx": "xl/workbook.xml",
            ".pptx": "ppt/presentation.xml",
        }[extension]
        if required not in names:
            raise _rejected("mime_extension_mismatch")

    def _inspect_image(self, data: bytes) -> str:
        normalized, mime = self.normalize_image(data=data)
        del normalized
        return mime

    def _extract_docx(self, data: bytes) -> ExtractedContent:
        from docx import Document

        document = Document(io.BytesIO(data))
        chunks = [p.text for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                chunks.append("\t".join(cell.text for cell in row.cells))
        return self._bounded("\n".join(chunks), [{"type": "document"}])

    def _extract_xlsx(self, data: bytes) -> ExtractedContent:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        chunks: list[str] = []
        segments: list[dict[str, object]] = []
        for sheet in workbook.worksheets:
            segments.append({"type": "worksheet", "name": sheet.title})
            chunks.append(f"[Sheet: {sheet.title}]")
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_index > self.settings.max_spreadsheet_rows:
                    break
                values = row[: self.settings.max_spreadsheet_columns]
                chunks.append("\t".join("" if value is None else str(value) for value in values))
        workbook.close()
        return self._bounded("\n".join(chunks), segments)

    def _extract_pptx(self, data: bytes) -> ExtractedContent:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        chunks: list[str] = []
        segments: list[dict[str, object]] = []
        for index, slide in enumerate(presentation.slides, start=1):
            if index > self.settings.max_slides:
                break
            segments.append({"type": "slide", "index": index})
            chunks.append(f"[Slide {index}]")
            for shape in slide.shapes:
                text = str(getattr(shape, "text", "")).strip()
                if text:
                    chunks.append(text)
        return self._bounded("\n".join(chunks), segments)

    def _bounded(self, text: str, segments: list[dict[str, object]]) -> ExtractedContent:
        normalized = text.replace("\x00", "").strip()
        truncated = len(normalized) > self.settings.max_extract_chars
        return ExtractedContent(
            text=normalized[: self.settings.max_extract_chars],
            segments=segments,
            parser_version=self.parser_version,
            truncated=truncated,
        )


def _rejected(code: str) -> NonRetryableExecutionError:
    return NonRetryableExecutionError(code, safe_message=code)
