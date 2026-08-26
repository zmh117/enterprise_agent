#!/usr/bin/env python3
"""Generate deterministic, non-business Docling benchmark inputs."""

from __future__ import annotations

import argparse
from io import BytesIO
from pathlib import Path

from docx import Document
from docx.shared import Inches as DocxInches
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches as PptxInches


def _image_bytes(image: Image.Image, image_format: str) -> bytes:
    output = BytesIO()
    image.save(output, format=image_format)
    return output.getvalue()


def _make_image() -> Image.Image:
    image = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((80, 90, 1520, 810), outline="black", width=6)
    draw.text((140, 180), "Synthetic Docling CPU benchmark", fill="black")
    draw.text((140, 280), "No customer or production data", fill="black")
    draw.text((140, 380), "Order 12345  Total 678.90", fill="black")
    return image


def _write_pdf(path: Path, *, pages: int, target_bytes: int | None = None) -> None:
    objects: dict[int, bytes] = {}
    page_ids = [4 + index * 2 for index in range(pages)]
    objects[1] = b"<< /Type /Catalog /Pages 2 0 R >>"
    objects[2] = (
        b"<< /Type /Pages /Count "
        + str(pages).encode("ascii")
        + b" /Kids ["
        + b" ".join(f"{item} 0 R".encode("ascii") for item in page_ids)
        + b"] >>"
    )
    objects[3] = b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>"
    for index, page_id in enumerate(page_ids, start=1):
        content_id = page_id + 1
        content = (
            b"BT /F1 12 Tf 72 720 Td (Synthetic benchmark page "
            + str(index).encode("ascii")
            + b") Tj ET\n"
        )
        objects[page_id] = (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 3 0 R >> >> /Contents "
            + f"{content_id} 0 R".encode("ascii")
            + b" >>"
        )
        objects[content_id] = (
            b"<< /Length "
            + str(len(content)).encode("ascii")
            + b" >>\nstream\n"
            + content
            + b"endstream"
        )

    body = bytearray(b"%PDF-1.4\n")
    offsets = [0] * (max(objects) + 1)
    for object_id in sorted(objects):
        offsets[object_id] = len(body)
        body.extend(f"{object_id} 0 obj\n".encode("ascii"))
        body.extend(objects[object_id])
        body.extend(b"\nendobj\n")
    xref_offset = len(body)
    body.extend(f"xref\n0 {len(offsets)}\n".encode("ascii"))
    body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        body.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    body.extend(
        b"trailer\n<< /Size "
        + str(len(offsets)).encode("ascii")
        + b" /Root 1 0 R >>\nstartxref\n"
        + str(xref_offset).encode("ascii")
        + b"\n%%EOF\n"
    )
    if target_bytes is not None:
        if target_bytes < len(body) + 2:
            raise ValueError("target PDF size is smaller than generated document")
        body.extend(b"%" + b"P" * (target_bytes - len(body) - 2) + b"\n")
    path.write_bytes(bytes(body))


def generate(output_dir: Path) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    image = _make_image()
    for suffix, image_format in (("png", "PNG"), ("jpg", "JPEG"), ("webp", "WEBP")):
        (output_dir / f"synthetic-image.{suffix}").write_bytes(
            _image_bytes(image, image_format)
        )

    document = Document()
    document.add_heading("Synthetic Docling benchmark", level=1)
    document.add_paragraph("No customer or production data is present in this file.")
    document.add_picture(
        BytesIO(_image_bytes(image, "PNG")), width=DocxInches(5.5)
    )
    table = document.add_table(rows=3, cols=2)
    for row_index, row in enumerate(table.rows):
        row.cells[0].text = f"Metric {row_index + 1}"
        row.cells[1].text = str((row_index + 1) * 10)
    document.save(output_dir / "synthetic-document.docx")

    presentation = Presentation()
    for slide_index in range(3):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title = slide.shapes.title
        if title is not None:
            title.text = f"Synthetic benchmark slide {slide_index + 1}"
        slide.shapes.add_picture(
            BytesIO(_image_bytes(image, "PNG")),
            PptxInches(1),
            PptxInches(1.5),
            width=PptxInches(8),
        )
    presentation.save(output_dir / "synthetic-presentation.pptx")

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Synthetic"
    worksheet.append(["row", "description", "amount"])
    for row_index in range(1, 101):
        worksheet.append([row_index, f"synthetic-{row_index}", row_index * 1.25])
    workbook.save(output_dir / "synthetic-workbook.xlsx")

    _write_pdf(output_dir / "synthetic-10-page.pdf", pages=10)
    _write_pdf(output_dir / "synthetic-100-page.pdf", pages=100)
    (output_dir / "synthetic-invalid.docx").write_bytes(
        b"synthetic invalid Office package; no customer data\n"
    )
    _write_pdf(
        output_dir / "synthetic-boundary-300-page.pdf",
        pages=300,
        target_bytes=25 * 1024 * 1024 - 1024,
    )


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("output_dir", type=Path)
    args = parser.parse_args()
    generate(args.output_dir)


if __name__ == "__main__":
    main()
