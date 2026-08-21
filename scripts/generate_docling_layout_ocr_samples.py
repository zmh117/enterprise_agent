#!/usr/bin/env python3
"""Generate deterministic synthetic Office samples for layout-OCR limits."""

from __future__ import annotations

import argparse
from io import BytesIO
from pathlib import Path

from docx import Document
from docx.shared import Inches as DocxInches
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches as PptxInches


IMAGE_WIDTH = 1280
IMAGE_HEIGHT = 720
TYPICAL_DOCX_PICTURES = 6
TYPICAL_PPTX_PICTURES = 8
BOUNDARY_PICTURES = 32


def _font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = (
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    )
    for candidate in candidates:
        path = Path(candidate)
        if path.is_file():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default(size=size)


def _picture(index: int) -> bytes:
    image = Image.new("RGB", (IMAGE_WIDTH, IMAGE_HEIGHT), "white")
    draw = ImageDraw.Draw(image)
    title_font = _font(44)
    body_font = _font(34)
    draw.rounded_rectangle((45, 45, 1235, 675), radius=20, outline="black", width=5)
    draw.text((90, 90), f"Synthetic layout image {index:02d}", fill="black", font=title_font)
    rows = (
        ("LEFT BLOCK", 100, 240),
        (f"ORDER {1000 + index}", 100, 360),
        ("RIGHT BLOCK", 760, 240),
        (f"TOTAL {index * 37}.50", 760, 360),
    )
    for value, x, y in rows:
        draw.rectangle((x - 20, y - 18, x + 430, y + 62), outline="#555555", width=3)
        draw.text((x, y), value, fill="black", font=body_font)
    output = BytesIO()
    image.save(output, format="PNG", optimize=True)
    return output.getvalue()


def _transform_picture() -> bytes:
    image = Image.new("RGB", (1200, 800), "white")
    draw = ImageDraw.Draw(image)
    font = _font(42)
    draw.rectangle((0, 0, 1200, 200), fill="#ff6666")
    draw.text((80, 70), "CROPPED TOP", fill="black", font=font)
    draw.rectangle((0, 200, 1200, 600), fill="#99dd99")
    draw.text((360, 360), "VISIBLE CENTER", fill="black", font=font)
    draw.rectangle((0, 600, 1200, 800), fill="#6699ff")
    draw.text((80, 670), "CROPPED BOTTOM", fill="black", font=font)
    output = BytesIO()
    image.save(output, format="PNG", optimize=True)
    return output.getvalue()


def _write_docx(path: Path, pictures: list[bytes]) -> None:
    document = Document()
    document.add_heading("Synthetic layout OCR benchmark", level=1)
    document.add_paragraph("No customer or production data is present.")
    for index, picture in enumerate(pictures, start=1):
        document.add_heading(f"Synthetic section {index}", level=2)
        document.add_paragraph(f"Stable parent text for picture {index}.")
        document.add_picture(BytesIO(picture), width=DocxInches(6.0))
    document.save(path)


def _write_pptx(path: Path, pictures: list[bytes]) -> None:
    presentation = Presentation()
    for index, picture in enumerate(pictures, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        if slide.shapes.title is not None:
            slide.shapes.title.text = f"Synthetic layout slide {index}"
        slide.shapes.add_picture(
            BytesIO(picture),
            PptxInches(1.0),
            PptxInches(1.5),
            width=PptxInches(8.0),
        )
    presentation.save(path)


def _write_transform_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    shape = slide.shapes.add_picture(
        BytesIO(_transform_picture()),
        PptxInches(2.0),
        PptxInches(1.0),
        width=PptxInches(6.0),
    )
    shape.crop_top = 0.25
    shape.crop_bottom = 0.25
    shape.rotation = 90
    presentation.save(path)


def generate(output_dir: Path) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    pictures = [_picture(index) for index in range(1, BOUNDARY_PICTURES + 1)]
    _write_docx(
        output_dir / "typical.docx",
        pictures[:TYPICAL_DOCX_PICTURES],
    )
    _write_pptx(
        output_dir / "typical.pptx",
        pictures[:TYPICAL_PPTX_PICTURES],
    )
    _write_docx(output_dir / "boundary.docx", pictures)
    _write_pptx(output_dir / "boundary.pptx", pictures)
    _write_transform_pptx(output_dir / "transform.pptx")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("output_dir", type=Path)
    args = parser.parse_args()
    generate(args.output_dir)


if __name__ == "__main__":
    main()
